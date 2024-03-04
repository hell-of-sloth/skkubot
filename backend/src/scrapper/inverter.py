import olefile # olefile 패키지 설치 필요
import os
import fitz # PyMuPDF 패키지 설치 필요
import docx # python-docx 패키지 설치 필요
from pptx import Presentation # python-pptx 패키지 설치 필요
import zlib # zlib 패키지 설치 필요
import struct


def fileInverter():
    
    # txt 파일 저장할 폴더 생성
    txt_path = os.path.dirname(os.path.abspath(__file__)) + "/txt_files"

    try:
        if not os.path.exists(txt_path):
            os.makedirs(txt_path) 
    except OSError:
        print('Error: Creating directory. ' + txt_path)
        
        
    # 한글, PDF, Word, PPT 파일을 txt 파일로 변환
    path = os.path.dirname(os.path.abspath(__file__)) 
    folder_path = path + "/download"

    for file in os.listdir(folder_path):
        try:
            # 한글 파일 .hwp
            if file.endswith(".hwp"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                f = olefile.OleFileIO(filename)
                dirs = f.listdir()
                
                # 문서 포맷 압축 여부 확인
                header = f.openstream("FileHeader")
                header_data = header.read()
                is_compressed = (header_data[36] & 1) == 1
                
                # Body section 추출
                nums = []
                for d in dirs:
                    if d[0] == "BodyText":
                        nums.append(int(d[1][len("Section"):]))
                sections = ["BodyText/Section" + str(i) for i in sorted(nums)]
                
                # text 추출
                
                hwp_text = ""
                for section in sections:
                    bodytext = f.openstream(section)
                    hwp_data = bodytext.read()
                    if is_compressed:
                        unpacked_data = zlib.decompress(hwp_data, -15)
                    else:
                        unpacked_data = hwp_data
                        
                    # 각 section의 text 추출
                    
                    section_text = ""
                    i = 0
                    size = len(unpacked_data)
                    while i < size:
                        header = struct.unpack_from("<l", unpacked_data, i)[0]
                        rec_type = header & 0x3ff
                        rec_len = (header >> 20) & 0xfff
                        
                        if rec_type in [67]:
                            rec_data = unpacked_data[i + 4:i + 4 + rec_len]
                            section_text += rec_data.decode("utf-16")
                            section_text += "\n"
                    
                        i += 4 + rec_len
                        
                    hwp_text += section_text
                    hwp_text += "\n"
                                    
                txt_f.write(hwp_text)
                txt_f.close()
                
            # PDF 파일 .pdf
            if file.endswith(".pdf"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                pdf = fitz.open(filename)
                for page in pdf:
                    text = page.get_text()
                    txt_f.write(text)
                txt_f.close()
                
            # Word 파일 .docx
            if file.endswith(".docx"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                doc = docx.Document(filename)
                for para in doc.paragraphs:
                    txt_f.write(para.text)
                txt_f.close()
                
            # PPT 파일 .pptx
            if file.endswith(".pptx"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                ppt = Presentation(filename)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt_f.write(shape.text)
                txt_f.close()
            
                
            print("Success: " + file)
        except:
            print("Error: " + file)

fileInverter()