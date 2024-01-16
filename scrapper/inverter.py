import olefile # olefile 패키지 설치 필요
import os
import fitz # PyMuPDF 패키지 설치 필요
import docx # python-docx 패키지 설치 필요
from pptx import Presentation # python-pptx 패키지 설치 필요
import subprocess


def fileInverter():
    
    # txt 파일 저장할 폴더 생성
    txt_path = os.path.dirname(os.path.abspath(__file__)) + "/txt_files"

    try:
        if not os.path.exists(txt_path):
            os.makedirs(txt_path) 
    except OSError:
        print('Error: Creating directory. ' + txt_path)
        
        
    # 한글, PDF, Word, PPT 파일을 txt 파일로 변환
    path = os.path.dirname(os.path.abspath(__file__)) 
    folder_path = path + "/download"

    for file in os.listdir(folder_path):
        try:
            # 한글 파일 .hwp
            if file.endswith(".hwp"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                f = olefile.OleFileIO(filename)
                encoded_text = f.openstream('PrvText').read()
                decoded_text = encoded_text.decode('UTF-16')
                txt_f.write(decoded_text)
                txt_f.close()
                
            # PDF 파일 .pdf
            if file.endswith(".pdf"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                pdf = fitz.open(filename)
                for page in pdf:
                    text = page.get_text()
                    txt_f.write(text)
                txt_f.close()
                
            # Word 파일 .docx
            if file.endswith(".docx"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                doc = docx.Document(filename)
                for para in doc.paragraphs:
                    txt_f.write(para.text)
                txt_f.close()
                
            # PPT 파일 .pptx
            if file.endswith(".pptx"):
                txt_f = open(txt_path + "/" + file + ".txt", 'w', encoding='UTF-8')
                filename = os.path.join(folder_path, file)
                ppt = Presentation(filename)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt_f.write(shape.text)
                txt_f.close()
            
                
            print("Success: " + file)
        except:
            print("Error: " + file)

fileInverter()